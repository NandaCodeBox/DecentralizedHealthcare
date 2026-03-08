#!/usr/bin/env python3
"""
Update the AWS AI Bharat Hackathon template PPTX with Arogya AI content.
Preserves the template background images and styling.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

def set_text_in_shape(shape, text, font_size=12, bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """Set text in a shape with formatting"""
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color

def add_bullet_text(shape, items, font_size=11, color=None):
    """Add bullet points to a shape"""
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = item
        run.font.size = Pt(font_size)
        if color:
            run.font.color.rgb = color

def update_presentation():
    print("=" * 80)
    print("UPDATING HACKATHON TEMPLATE PPTX")
    print("=" * 80)

    template_path = r"Deck\AWS_AI_Bharat_Hackathon_Final_Submission.pptx"
    output_path = r"Deck\Arogya_AI_Final_Submission.pptx"

    if not os.path.exists(template_path):
        print(f"✗ Template not found: {template_path}")
        return False

    prs = Presentation(template_path)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    DARK = RGBColor(0x33, 0x33, 0x33)
    BLUE = RGBColor(0x00, 0x6B, 0xBF)

    print("\n✓ Template loaded (14 slides)")

    # ============================================================
    # SLIDE 1: Team Info (index 0)
    # Shape 1 = "Team Name:", Shape 3 = "Problem Statement:", Shape 2 = "Team Leader:"
    # ============================================================
    print("  Slide 1: Team Info...")
    slide = prs.slides[0]
    shapes = list(slide.shapes)
    # shapes[0] = background image, shapes[1..3] = text boxes
    set_text_in_shape(shapes[1], "Team Name : NandaCodeBox", 18, True, WHITE)
    set_text_in_shape(shapes[3], "Problem Statement : AI-Powered Healthcare Access for 900M Underserved Indians", 14, True, WHITE)
    set_text_in_shape(shapes[2], "Team Leader Name : Nanda Kumar N", 14, True, WHITE)

    # ============================================================
    # SLIDE 2: Brief about the Idea (index 1)
    # ============================================================
    print("  Slide 2: Brief about the Idea...")
    slide = prs.slides[1]
    shapes = list(slide.shapes)
    set_text_in_shape(shapes[1], "Brief about the Idea:", 16, True, WHITE)
    # Add content text box
    from pptx.util import Inches
    txBox = slide.shapes.add_textbox(Inches(0.4), Inches(1.8), Inches(9.0), Inches(3.2))
    add_bullet_text(txBox, [
        "Arogya.AI — AI-powered healthcare triage platform for 900M underserved Indians",
        "Patients describe symptoms in 10 Indian languages via simple mobile interface",
        "Amazon Bedrock AI analyzes symptoms in 3 seconds with 94% confidence",
        "3 Autonomous Agentic AI agents (Supervisor Validation, Care Pathway, Clinical Decision)",
        "6-level multi-reasoning: Symptom → Vitals → Urgency → Facility → Care Plan → Risk",
        "81% cases auto-approved without human intervention (38/47 today)",
        "Smart facility matching (95% AI match) with instant appointment booking",
        "30 seconds vs 45 minutes traditional triage — 90X faster",
        "Cost: $6.22/month for 10,000 patients ($0.0006 per patient)",
        "Live on AWS: S3 + Lambda + Bedrock + DynamoDB — fully serverless"
    ], 11, WHITE)

    # ============================================================
    # SLIDE 3: Why AI / AWS / Value (index 2)
    # ============================================================
    print("  Slide 3: Why AI + AWS + Value...")
    slide = prs.slides[2]
    shapes = list(slide.shapes)
    tf = shapes[1].text_frame
    tf.clear()
    tf.word_wrap = True

    sections = [
        ("Why AI is required:", [
            "• 900M Indians lack timely healthcare access — human triage takes 45 min",
            "• AI enables instant symptom analysis with 94% confidence in 3 seconds",
            "• 3 Agentic AI agents autonomously handle 81% of cases 24/7",
            "• 6-level reasoning mimics expert clinical decision-making",
        ]),
        ("How AWS services are used:", [
            "• Amazon Bedrock — AI reasoning engine for all 3 agents",
            "• AWS Lambda — Serverless compute for agent execution",
            "• Amazon S3 — Static website hosting for React frontend",
            "• Amazon DynamoDB — Patient data and case storage",
            "• Amazon Polly — Multilingual voice support (10 languages)",
            "• CloudWatch — Monitoring and observability",
        ]),
        ("Value AI adds to user experience:", [
            "• Instant triage: 30s vs 45 min (90X faster)",
            "• Multilingual: Healthcare in your mother tongue (10 languages)",
            "• Smart matching: 95% AI facility match accuracy",
            "• Autonomous: 81% cases need zero human intervention",
            "• Affordable: $0.0006 per patient (0.06 cents)",
        ]),
    ]

    for i, (title, bullets) in enumerate(sections):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            p.space_before = Pt(8)
        run = p.add_run()
        run.text = title
        run.font.size = Pt(12)
        run.font.bold = True
        run.font.color.rgb = WHITE

        for bullet in bullets:
            p = tf.add_paragraph()
            p.space_before = Pt(2)
            run = p.add_run()
            run.text = bullet
            run.font.size = Pt(9)
            run.font.color.rgb = WHITE

    # ============================================================
    # SLIDE 4: Features (index 3)
    # ============================================================
    print("  Slide 4: Features...")
    slide = prs.slides[3]
    shapes = list(slide.shapes)
    tf = shapes[1].text_frame
    tf.clear()
    tf.word_wrap = True

    features = [
        "List of features offered by the solution:",
        "",
        "🏥 Patient Features:",
        "  • Multilingual symptom intake (10 Indian languages)",
        "  • Simple tap-based symptom tiles (no medical jargon)",
        "  • AI-powered triage with 94% confidence score",
        "  • Smart facility recommendations (95% AI match)",
        "  • Instant appointment booking with confirmation",
        "  • Emergency detection and priority routing",
        "",
        "🤖 Agentic AI Features:",
        "  • 3 Autonomous AI Agents (Supervisor, Care Pathway, Clinical Decision)",
        "  • 6-Level Multi-Reasoning Process",
        "  • 81% Auto-Approval Rate (38/47 cases)",
        "  • Intelligent Escalation (knows when to ask humans)",
        "  • 24/7 Operation — no downtime",
        "",
        "👨‍⚕️ Supervisor Features:",
        "  • Real-time dashboard with case management",
        "  • Agentic AI toggle (purple indicator)",
        "  • Case review with AI reasoning transparency",
        "  • Override capability for escalated cases",
        "  • Performance analytics and metrics",
    ]

    for i, line in enumerate(features):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        if i == 0 or line.startswith("🏥") or line.startswith("🤖") or line.startswith("👨"):
            run.font.size = Pt(11)
            run.font.bold = True
        else:
            run.font.size = Pt(9)
        run.font.color.rgb = WHITE

    # ============================================================
    # SLIDE 5: Process Flow (index 4)
    # ============================================================
    print("  Slide 5: Process Flow...")
    slide = prs.slides[4]
    shapes = list(slide.shapes)
    tf = shapes[1].text_frame
    tf.clear()
    tf.word_wrap = True

    flow = [
        "Process Flow — Patient to Resolution:",
        "",
        "Patient Journey (Mobile View):",
        "  1. Open App → Select Language (10 options)",
        "  2. Login → Patient Dashboard",
        "  3. Tap Symptom Tiles → Fill Details → Rate Severity",
        "  4. Submit → AI Analysis (3 seconds)",
        "  5. View Triage Results (94% confidence)",
        "  6. See Facility Recommendations (95% match)",
        "  7. Book Appointment → Confirmation",
        "",
        "Agentic AI Pipeline (Behind the Scenes):",
        "  1. Supervisor Validation Agent → 6-level reasoning",
        "  2. Care Pathway Agent → Treatment plan + scheduling",
        "  3. Clinical Decision Agent → Diagnosis + recommendations",
        "  4. Auto-Approve (81%) or Escalate to Human (19%)",
        "",
        "Supervisor Review (Desktop View):",
        "  1. Dashboard → View all cases + statistics",
        "  2. Green = Auto-Approved | Orange = Needs Review",
        "  3. Review AI reasoning → Approve/Override",
    ]

    for i, line in enumerate(flow):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        if ":" in line and not line.startswith("  "):
            run.font.size = Pt(11)
            run.font.bold = True
        else:
            run.font.size = Pt(9)
        run.font.color.rgb = WHITE

    # ============================================================
    # SLIDE 6: Wireframes/Mockups (index 5)
    # ============================================================
    print("  Slide 6: Wireframes...")
    slide = prs.slides[5]
    shapes = list(slide.shapes)
    tf = shapes[1].text_frame
    tf.clear()
    tf.word_wrap = True

    wireframes = [
        "Wireframes / UI Screens:",
        "",
        "📱 Mobile Patient App (390x844):",
        "  • Homepage with emergency banner",
        "  • Language selector (10 Indian languages)",
        "  • Symptom tile grid (tap-based selection)",
        "  • Severity rating + duration input",
        "  • AI triage results with confidence score",
        "  • Facility cards with AI match percentage",
        "  • Appointment booking form",
        "",
        "💻 Desktop Supervisor Dashboard (1920x1080):",
        "  • Case management table with status indicators",
        "  • Agentic AI toggle (purple)",
        "  • Statistics panel (cases, auto-approved, rate)",
        "  • Case detail view with 6-level reasoning",
        "  • Green/Orange status indicators",
        "  • AI-powered provider search",
        "",
        "🌐 Live URL:",
        "  http://arogya-ai-healthcare-20260308102925",
        "  .s3-website-us-east-1.amazonaws.com",
    ]

    for i, line in enumerate(wireframes):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        if line.startswith("📱") or line.startswith("💻") or line.startswith("🌐") or i == 0:
            run.font.size = Pt(11)
            run.font.bold = True
        else:
            run.font.size = Pt(9)
        run.font.color.rgb = WHITE

    # ============================================================
    # SLIDE 7: Architecture (index 6)
    # ============================================================
    print("  Slide 7: Architecture...")
    slide = prs.slides[6]
    shapes = list(slide.shapes)
    tf = shapes[1].text_frame
    tf.clear()
    tf.word_wrap = True

    arch = [
        "Architecture — AWS Serverless + Agentic AI:",
        "",
        "Frontend Layer:",
        "  • React 18 + TypeScript + Tailwind CSS",
        "  • Hosted on Amazon S3 (Static Website)",
        "  • Mobile-first responsive design",
        "",
        "AI Agent Layer (3 Autonomous Agents):",
        "  • Agent 1: Supervisor Validation (AWS Lambda)",
        "    → 6-level reasoning, auto-approve 81%",
        "  • Agent 2: Care Pathway Orchestrator (AWS Lambda)",
        "    → Treatment plans, scheduling, coordination",
        "  • Agent 3: Clinical Decision Support (AWS Lambda)",
        "    → Differential diagnosis, recommendations",
        "",
        "AWS Services:",
        "  • Amazon Bedrock — AI reasoning engine",
        "  • AWS Lambda — Serverless agent execution",
        "  • Amazon DynamoDB — Patient data store",
        "  • Amazon S3 — Frontend hosting",
        "  • Amazon Polly — Multilingual voice",
        "  • CloudWatch — Monitoring & logging",
        "  • IAM — Security & access control",
    ]

    for i, line in enumerate(arch):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        if ":" in line and not line.startswith("  "):
            run.font.size = Pt(11)
            run.font.bold = True
        else:
            run.font.size = Pt(9)
        run.font.color.rgb = WHITE

    # ============================================================
    # SLIDE 8: Technologies (index 7)
    # ============================================================
    print("  Slide 8: Technologies...")
    slide = prs.slides[7]
    shapes = list(slide.shapes)
    tf = shapes[1].text_frame
    tf.clear()
    tf.word_wrap = True

    tech = [
        "Technologies Utilized:",
        "",
        "☁️ AWS Cloud Services:",
        "  • Amazon Bedrock (Claude/Titan) — AI reasoning",
        "  • AWS Lambda (Python 3.12) — Agent execution",
        "  • Amazon S3 — Static website hosting",
        "  • Amazon DynamoDB — NoSQL database",
        "  • Amazon Polly — Text-to-speech (10 languages)",
        "  • Amazon CloudWatch — Monitoring",
        "  • AWS IAM — Security & access control",
        "",
        "🖥️ Frontend:",
        "  • React 18 + TypeScript",
        "  • Tailwind CSS + Shadcn/UI",
        "  • i18next (internationalization)",
        "",
        "🤖 AI/ML:",
        "  • Amazon Bedrock (foundation models)",
        "  • 6-level multi-reasoning pipeline",
        "  • 3 autonomous agentic AI agents",
        "",
        "🧪 Testing:",
        "  • Playwright (automated E2E testing)",
        "  • Multi-language test suite (5 languages)",
    ]

    for i, line in enumerate(tech):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        if line.startswith("☁") or line.startswith("🖥") or line.startswith("🤖") or line.startswith("🧪") or i == 0:
            run.font.size = Pt(11)
            run.font.bold = True
        else:
            run.font.size = Pt(9)
        run.font.color.rgb = WHITE

    # ============================================================
    # SLIDE 9: Cost (index 8)
    # ============================================================
    print("  Slide 9: Implementation Cost...")
    slide = prs.slides[8]
    shapes = list(slide.shapes)
    tf = shapes[1].text_frame
    tf.clear()
    tf.word_wrap = True

    cost = [
        "Estimated Implementation Cost:",
        "",
        "Monthly Cost Breakdown (10,000 patients/month):",
        "  • AWS Lambda (3 agents): $0.62/month",
        "  • Amazon Bedrock AI: $3.50/month",
        "  • Amazon DynamoDB: $1.25/month",
        "  • Amazon S3 Hosting: $0.50/month",
        "  • Amazon Polly: $0.25/month",
        "  • CloudWatch: $0.10/month",
        "  ─────────────────────────────",
        "  Total: $6.22/month",
        "",
        "Per Patient Cost:",
        "  • $6.22 ÷ 10,000 = $0.000622 per patient",
        "  • That's 0.06 cents (six hundredths of one cent)",
        "",
        "ROI Analysis:",
        "  • Traditional triage cost: ~$15 per patient",
        "  • Arogya AI cost: $0.0006 per patient",
        "  • Savings: $14.999 per patient",
        "  • ROI: 241,000%",
        "",
        "Scaling:",
        "  • 100K patients: ~$62/month",
        "  • 1M patients: ~$620/month",
        "  • 10M patients: ~$6,200/month",
    ]

    for i, line in enumerate(cost):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        if ":" in line and not line.startswith("  "):
            run.font.size = Pt(11)
            run.font.bold = True
        else:
            run.font.size = Pt(9)
        run.font.color.rgb = WHITE

    # ============================================================
    # SLIDE 10: Snapshots (index 9)
    # ============================================================
    print("  Slide 10: Prototype Snapshots...")
    slide = prs.slides[9]
    shapes = list(slide.shapes)
    tf = shapes[1].text_frame
    tf.clear()
    tf.word_wrap = True

    snaps = [
        "Snapshots of the Prototype:",
        "",
        "📱 Patient Mobile App:",
        "  • Homepage with emergency banner and language selector",
        "  • Symptom tile grid — tap Chest Pain, Fever, etc.",
        "  • AI Triage Results — 94% confidence, High Priority",
        "  • Facility Cards — 95% AI match, 2km away",
        "  • Appointment Booking — date, time, confirmation",
        "",
        "💻 Supervisor Desktop Dashboard:",
        "  • Case management table with 47 cases",
        "  • Purple Agentic AI toggle (always on)",
        "  • Statistics: 38/47 auto-approved (81%)",
        "  • Case detail: 6-level reasoning breakdown",
        "  • Green indicator (auto-approved)",
        "  • Orange indicator (escalated to human)",
        "",
        "🌐 Live Application:",
        "  http://arogya-ai-healthcare-20260308102925",
        "  .s3-website-us-east-1.amazonaws.com",
        "  Patient: patient@arogya.ai / PatientPass123!",
        "  Supervisor: supervisor@arogya.ai / SupervisorPass123!",
    ]

    for i, line in enumerate(snaps):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        if line.startswith("📱") or line.startswith("💻") or line.startswith("🌐") or i == 0:
            run.font.size = Pt(11)
            run.font.bold = True
        else:
            run.font.size = Pt(9)
        run.font.color.rgb = WHITE

    # ============================================================
    # SLIDE 11: Performance (index 10)
    # ============================================================
    print("  Slide 11: Performance Benchmarking...")
    slide = prs.slides[10]
    shapes = list(slide.shapes)
    tf = shapes[1].text_frame
    tf.clear()
    tf.word_wrap = True

    perf = [
        "Prototype Performance Report / Benchmarking:",
        "",
        "⚡ Speed:",
        "  • AI Triage: 3 seconds (vs 45 min traditional)",
        "  • 90X faster than manual triage",
        "  • End-to-end patient journey: 30 seconds",
        "",
        "🎯 Accuracy:",
        "  • AI Confidence Score: 94%",
        "  • Facility Match Accuracy: 95%",
        "  • 6-level reasoning depth",
        "",
        "🤖 Automation:",
        "  • Auto-approval rate: 81% (38/47 cases)",
        "  • Intelligent escalation: 19% (9/47 cases)",
        "  • 24/7 availability — zero downtime",
        "",
        "🌍 Scale:",
        "  • 10 Indian languages supported",
        "  • Serverless — auto-scales to millions",
        "  • $0.0006 per patient cost",
        "",
        "🔒 Reliability:",
        "  • AWS Lambda 99.95% SLA",
        "  • DynamoDB 99.999% availability",
        "  • CloudWatch monitoring active",
    ]

    for i, line in enumerate(perf):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        if line.startswith("⚡") or line.startswith("🎯") or line.startswith("🤖") or line.startswith("🌍") or line.startswith("🔒") or i == 0:
            run.font.size = Pt(11)
            run.font.bold = True
        else:
            run.font.size = Pt(9)
        run.font.color.rgb = WHITE

    # ============================================================
    # SLIDE 12: Future Development (index 11)
    # ============================================================
    print("  Slide 12: Future Development...")
    slide = prs.slides[11]
    shapes = list(slide.shapes)
    tf = shapes[1].text_frame
    tf.clear()
    tf.word_wrap = True

    future = [
        "Additional Details / Future Development:",
        "",
        "🚀 Phase 2 (Next 3 months):",
        "  • Integration with ABDM (Ayushman Bharat Digital Mission)",
        "  • Real-time vitals monitoring via IoT devices",
        "  • Telemedicine video consultation integration",
        "  • Prescription management and pharmacy connect",
        "",
        "🌐 Phase 3 (6-12 months):",
        "  • Expand to 22 Indian languages",
        "  • Partner with PHCs (Primary Health Centers)",
        "  • Government health scheme integration",
        "  • Offline-first mode for low connectivity areas",
        "",
        "🤖 AI Enhancements:",
        "  • Continuous learning from case outcomes",
        "  • Predictive health alerts",
        "  • Medical image analysis (X-ray, ECG)",
        "  • Voice-based symptom input",
        "",
        "📊 Impact Target:",
        "  • Year 1: 1M patients served",
        "  • Year 3: 50M patients served",
        "  • Year 5: 200M patients — covering rural India",
    ]

    for i, line in enumerate(future):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        if line.startswith("🚀") or line.startswith("🌐") or line.startswith("🤖") or line.startswith("📊") or i == 0:
            run.font.size = Pt(11)
            run.font.bold = True
        else:
            run.font.size = Pt(9)
        run.font.color.rgb = WHITE

    # ============================================================
    # SLIDE 13: Prototype Assets (index 12)
    # ============================================================
    print("  Slide 13: Prototype Assets...")
    slide = prs.slides[12]
    shapes = list(slide.shapes)
    tf = shapes[1].text_frame
    tf.clear()
    tf.word_wrap = True

    assets = [
        "Prototype Assets:",
        "",
        "📦 GitHub Public Repository:",
        "  https://github.com/NandaCodeBox/DecentralizedHealthcare",
        "",
        "🎬 Demo Video (3 Minutes):",
        "  Video/Arogya_AI_Final_Synced.mp4",
        "  • Patient journey in mobile view (390x844)",
        "  • Supervisor dashboard in desktop view (1920x1080)",
        "  • Professional voiceover with story-telling",
        "  • All features + Agentic AI demonstrated",
        "",
        "🌐 Live Application:",
        "  http://arogya-ai-healthcare-20260308102925",
        "  .s3-website-us-east-1.amazonaws.com",
        "",
        "🔑 Credentials:",
        "  Patient: patient@arogya.ai / PatientPass123!",
        "  Supervisor: supervisor@arogya.ai / SupervisorPass123!",
        "",
        "🤖 Agent Endpoints (AWS Lambda):",
        "  Supervisor Agent:",
        "    https://35v66sz7u43rqq67e5fqmh6yeu0svwme.lambda-url.us-east-1.on.aws/",
        "  Care Pathway Agent:",
        "    https://kfrboux5jjxxtteqkp44e3psca0nzcic.lambda-url.us-east-1.on.aws/",
        "  Clinical Decision Agent:",
        "    https://46cvklukkhccawngp5g2yd7fpi0tqswa.lambda-url.us-east-1.on.aws/",
    ]

    for i, line in enumerate(assets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        if line.startswith("📦") or line.startswith("🎬") or line.startswith("🌐") or line.startswith("🔑") or line.startswith("🤖") or i == 0:
            run.font.size = Pt(11)
            run.font.bold = True
        else:
            run.font.size = Pt(9)
        run.font.color.rgb = WHITE

    # ============================================================
    # SAVE
    # ============================================================
    print(f"\n  Saving to: {output_path}")
    prs.save(output_path)
    
    print(f"\n✓ Presentation saved: {output_path}")
    print(f"  Slides: {len(prs.slides)}")
    return True


if __name__ == '__main__':
    print("\nUpdating hackathon template with Arogya AI content...\n")
    success = update_presentation()
    if success:
        print("\n" + "=" * 80)
        print("SUCCESS! Presentation ready: Deck/Arogya_AI_Final_Submission.pptx")
        print("=" * 80)
    else:
        print("\nFAILED — check errors above")
