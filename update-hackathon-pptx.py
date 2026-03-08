#!/usr/bin/env python3
"""
Update Hackathon PowerPoint Presentation
Adds latest demo video information and Agentic AI details
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def update_presentation():
    print("=" * 80)
    print("UPDATING HACKATHON POWERPOINT PRESENTATION")
    print("=" * 80)
    
    # Open the latest presentation
    pptx_file = "Deck/Arogya_AI_Hackathon_Final_Presentation.pptx"
    
    if not os.path.exists(pptx_file):
        print(f"\n✗ PowerPoint file not found: {pptx_file}")
        return False
    
    print(f"\n✓ Opening: {pptx_file}")
    
    try:
        prs = Presentation(pptx_file)
        print(f"  Current slides: {len(prs.slides)}")
        
        # Add new slide: Demo Video Information
        print("\n→ Adding Demo Video slide...")
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = "📹 Demo Video - Story-Driven Presentation"
        
        # Content
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        # Add content
        p = tf.paragraphs[0]
        p.text = "3-Minute Hackathon Demo"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 128, 128)
        
        # Segment 1
        p = tf.add_paragraph()
        p.text = "Segment 1: The Problem (0-30s)"
        p.level = 0
        p.font.size = Pt(18)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = "Rajesh's story - 900M Indians in underserved areas"
        p.level = 1
        p.font.size = Pt(16)
        
        # Segment 2
        p = tf.add_paragraph()
        p.text = "Segment 2: Patient Journey (30-75s)"
        p.level = 0
        p.font.size = Pt(18)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = "Mobile app → Symptom tiles → AI triage → Appointment booking"
        p.level = 1
        p.font.size = Pt(16)
        
        # Segment 3
        p = tf.add_paragraph()
        p.text = "Segment 3: Agentic AI (75-135s)"
        p.level = 0
        p.font.size = Pt(18)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = "Supervisor dashboard → 6-level reasoning → 81% automation"
        p.level = 1
        p.font.size = Pt(16)
        
        # Segment 4
        p = tf.add_paragraph()
        p.text = "Segment 4: Multi-language & Scale (135-165s)"
        p.level = 0
        p.font.size = Pt(18)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = "10 languages → AI provider search → AWS deployment"
        p.level = 1
        p.font.size = Pt(16)
        
        # Segment 5
        p = tf.add_paragraph()
        p.text = "Segment 5: The Impact (165-180s)"
        p.level = 0
        p.font.size = Pt(18)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = "90X faster → $0.0006/patient → Democratizing healthcare"
        p.level = 1
        p.font.size = Pt(16)
        
        # Add new slide: Agentic AI Details
        print("→ Adding Agentic AI Details slide...")
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "🤖 Agentic AI System - Three Autonomous Agents"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        # Agent 1
        p = tf.paragraphs[0]
        p.text = "Agent 1: Supervisor Validation"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(128, 0, 128)
        
        p = tf.add_paragraph()
        p.text = "6-level multi-reasoning → Auto-approves 70-80% of cases"
        p.level = 1
        p.font.size = Pt(16)
        
        p = tf.add_paragraph()
        p.text = "Endpoint: https://35v66sz7u43rqq67e5fqmh6yeu0svwme.lambda-url.us-east-1.on.aws/"
        p.level = 1
        p.font.size = Pt(12)
        
        # Agent 2
        p = tf.add_paragraph()
        p.text = "Agent 2: Care Pathway Orchestrator"
        p.level = 0
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(128, 0, 128)
        
        p = tf.add_paragraph()
        p.text = "Autonomous care coordination → Auto-scheduling"
        p.level = 1
        p.font.size = Pt(16)
        
        # Agent 3
        p = tf.add_paragraph()
        p.text = "Agent 3: Clinical Decision Support"
        p.level = 0
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(128, 0, 128)
        
        p = tf.add_paragraph()
        p.text = "Differential diagnosis → Treatment recommendations"
        p.level = 1
        p.font.size = Pt(16)
        
        # Cost
        p = tf.add_paragraph()
        p.text = "💰 Cost: $6.22/month for 10,000 patients"
        p.level = 0
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 128, 0)
        
        # Add new slide: Technical Architecture
        print("→ Adding Technical Architecture slide...")
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "🏗️ Technical Architecture - AWS Serverless"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        p = tf.paragraphs[0]
        p.text = "Frontend"
        p.font.size = Pt(20)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = "Next.js + React + TypeScript → S3 + CloudFront"
        p.level = 1
        p.font.size = Pt(16)
        
        p = tf.add_paragraph()
        p.text = "Live: http://arogya-ai-healthcare-20260308102925.s3-website-us-east-1.amazonaws.com"
        p.level = 1
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "Backend - Agentic AI"
        p.level = 0
        p.font.size = Pt(20)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = "3 AWS Lambda Functions (Python 3.11)"
        p.level = 1
        p.font.size = Pt(16)
        
        p = tf.add_paragraph()
        p.text = "Amazon Bedrock (Claude 3 Haiku) for AI reasoning"
        p.level = 1
        p.font.size = Pt(16)
        
        p = tf.add_paragraph()
        p.text = "DynamoDB for state management"
        p.level = 1
        p.font.size = Pt(16)
        
        p = tf.add_paragraph()
        p.text = "SNS for notifications"
        p.level = 1
        p.font.size = Pt(16)
        
        p = tf.add_paragraph()
        p.text = "Key Metrics"
        p.level = 0
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 128, 128)
        
        p = tf.add_paragraph()
        p.text = "⚡ 30-second triage (90X faster than 45 minutes)"
        p.level = 1
        p.font.size = Pt(16)
        
        p = tf.add_paragraph()
        p.text = "🤖 81% automation rate (38/47 cases auto-approved)"
        p.level = 1
        p.font.size = Pt(16)
        
        p = tf.add_paragraph()
        p.text = "💰 $0.000622 per patient (223,000% ROI)"
        p.level = 1
        p.font.size = Pt(16)
        
        # Add new slide: Impact & ROI
        print("→ Adding Impact & ROI slide...")
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "📊 Impact & ROI - Transforming Healthcare"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        p = tf.paragraphs[0]
        p.text = "Current Scale (10,000 patients/month)"
        p.font.size = Pt(20)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = "AI System Cost: $6.22/month"
        p.level = 1
        p.font.size = Pt(16)
        
        p = tf.add_paragraph()
        p.text = "Traditional System: $15,000/month (staff salaries)"
        p.level = 1
        p.font.size = Pt(16)
        
        p = tf.add_paragraph()
        p.text = "Monthly Savings: $14,993.78"
        p.level = 1
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(0, 128, 0)
        
        p = tf.add_paragraph()
        p.text = "ROI: 241,000%"
        p.level = 1
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 128, 0)
        
        p = tf.add_paragraph()
        p.text = "India Scale (1,000,000 patients/month)"
        p.level = 0
        p.font.size = Pt(20)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = "AI System Cost: $681/month"
        p.level = 1
        p.font.size = Pt(16)
        
        p = tf.add_paragraph()
        p.text = "Traditional System: $1,500,000/month"
        p.level = 1
        p.font.size = Pt(16)
        
        p = tf.add_paragraph()
        p.text = "Monthly Savings: $1,499,319"
        p.level = 1
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(0, 128, 0)
        
        p = tf.add_paragraph()
        p.text = "ROI: 220,000%"
        p.level = 1
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 128, 0)
        
        p = tf.add_paragraph()
        p.text = "🎯 Target: 900M Indians in underserved areas"
        p.level = 0
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(128, 0, 128)
        
        # Save updated presentation
        output_file = "Deck/Arogya_AI_Hackathon_Updated.pptx"
        prs.save(output_file)
        
        print(f"\n✓ Presentation updated: {output_file}")
        print(f"  Total slides: {len(prs.slides)}")
        print(f"  New slides added: 4")
        
        return True
        
    except Exception as e:
        print(f"\n✗ Error updating presentation: {e}")
        return False

if __name__ == '__main__':
    print("\nUpdating PowerPoint presentation with latest demo and Agentic AI info...\n")
    
    success = update_presentation()
    
    if success:
        print("\n" + "=" * 80)
        print("SUCCESS! 🎉")
        print("=" * 80)
        print("\nPowerPoint presentation updated with:")
        print("  ✓ Demo video segments and timing")
        print("  ✓ Agentic AI system details")
        print("  ✓ Technical architecture")
        print("  ✓ Impact and ROI metrics")
        print("\nFile: Deck/Arogya_AI_Hackathon_Updated.pptx")
    else:
        print("\n" + "=" * 80)
        print("FAILED")
        print("=" * 80)
        print("\nPlease check the error messages above.")
